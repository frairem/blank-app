from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import pandas as pd
import re

def apply_bold_to_paragraph(paragraph):
    """
    Applies bold formatting to words enclosed between ** iterating over each part of each pptx text_frame.paragraphs.

    Example:
    paragraph = 'Responsible for **collecting** and **analyzing** data.'
    parts = ['Responsible for ','**collecting**',' and' ,'**analyzing**',' data']
    part = '**analyzing**'
    result = analyzing in bold
    """
    text = paragraph.text
    parts = re.split(r'(\*\*.*?\*\*)', text)

    # Remove existing runs
    for _ in range(len(paragraph.runs)):
        paragraph._element.remove(paragraph.runs[0]._r)
        
    # Rebuild runs with bold formatting
    for part in parts:
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]  # remove ** markers
            run.font.bold = True
        else:
            run.text = part

def populate_pptx(df_text):
    """
    Populates a PowerPoint presentation with text from df_text.
    Each row of df_text should contain 'section_name' and 'output' columns.
    """
    prs = Presentation('data/input/SC&O OP VACIO - TEMPLATE 3 - New.pptx')
    slide = prs.slides[0]

    # Define sections where bullets should be applied
    bullet_sections = {"industry experience", "functional experience", "certifications/training"}

    for index, row in df_text.iterrows():
        section_name = row["section_name"].strip().lower()
        new_text = str(row["output"]).strip()
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if shape.name.strip().lower() != row['section_name'].strip().lower():
                continue

            text_frame = shape.text_frame
            text_frame.clear()
            # new_text = row['output']

            # --- Format lines depending on section ---
            lines = new_text.splitlines()
            if section_name.startswith("role"):
                # Keep title without bullet and make it bold
                title_line = f"**{lines[0].strip()}**"
                formatted_lines = "\n".join([title_line] + [f"• {line}" for line in lines[1:]])
            elif section_name in bullet_sections:
                formatted_lines = "\n".join([f"• {line}" for line in lines])
            else:
                # No bullets for other sections
                formatted_lines = "\n".join(lines)

            text_frame.text = formatted_lines

            # Apply styles
            for paragraph in text_frame.paragraphs:
                apply_bold_to_paragraph(paragraph)

                # Handle NAME formatting
                if shape.name.strip().lower() == 'name':
                    run = paragraph.runs[0]
                    run.font.name = "Graphik Black"
                    run.font.size = Pt(42)
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.text = run.text.upper().strip()

                # Handle TOWER formatting
                elif shape.name.strip().lower() == 'tower':
                    run = paragraph.runs[0]
                    run.font.name = "Graphik Black"
                    run.font.size = Pt(24)
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.text = run.text.upper().strip()

                # Default formatting for other sections
                else:
                    for run in paragraph.runs:
                        run.font.name = "Graphik"
                        run.font.size = Pt(9)
                        run.font.color.rgb = RGBColor(0, 0, 0)

    output_path = "data/output/updated_presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation was successfully created")
    return output_path